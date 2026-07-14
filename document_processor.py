import os
import fitz  # PyMuPDF
from pptx import Presentation
from langchain_text_splitters import RecursiveCharacterTextSplitter

class DocumentProcessor:
    def __init__(self, data_dir="data", chunk_size=1000, chunk_overlap=50):
        self.data_dir = data_dir
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
            length_function=len
        )

    def extract_text_from_pdf(self, file_path):
        documents = []
        filename = os.path.basename(file_path)
        try:
            doc = fitz.open(file_path)
            for page_num in range(len(doc)):
                page = doc[page_num]
                text = page.get_text()
                if text.strip():
                    documents.append({
                        "text": text,
                        "metadata": {"source": filename, "page": page_num + 1}
                    })
            doc.close()
        except Exception as e:
            print(f"Error reading PDF {filename}: {e}")
        return documents

    def extract_text_from_pptx(self, file_path):
        documents = []
        filename = os.path.basename(file_path)
        try:
            prs = Presentation(file_path)
            for slide_num, slide in enumerate(prs.slides):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                
                text = "\n".join(slide_text)
                if text.strip():
                    documents.append({
                        "text": text,
                        "metadata": {"source": filename, "page": slide_num + 1}
                    })
        except Exception as e:
            print(f"Error reading PPTX {filename}: {e}")
        return documents

    def extract_text_from_txt(self, file_path):
        documents = []
        filename = os.path.basename(file_path)
        try:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as file:
                text = file.read()
            if text.strip():
                documents.append({
                    "text": text,
                    "metadata": {"source": filename, "page": 1}
                })
        except Exception as e:
            print(f"Error reading text file {filename}: {e}")
        return documents

    def process_file(self, file_path):
        filename = os.path.basename(file_path)
        lower_name = filename.lower()

        if lower_name.endswith(".pdf"):
            raw_documents = self.extract_text_from_pdf(file_path)
        elif lower_name.endswith(".pptx"):
            raw_documents = self.extract_text_from_pptx(file_path)
        elif lower_name.endswith(".txt"):
            raw_documents = self.extract_text_from_txt(file_path)
        else:
            return []

        final_chunks = []
        for doc in raw_documents:
            chunks = self.text_splitter.split_text(doc["text"])
            for chunk_index, chunk in enumerate(chunks, start=1):
                metadata = dict(doc["metadata"])
                metadata["chunk"] = chunk_index
                final_chunks.append({
                    "text": chunk,
                    "metadata": metadata
                })

        return final_chunks

    def process_all_documents(self):
        if not os.path.exists(self.data_dir):
            os.makedirs(self.data_dir)
            print(f"Created empty '{self.data_dir}' folder. Please drop your PDFs/PPTXs inside it.")
            return []

        final_chunks = []
        for filename in os.listdir(self.data_dir):
            file_path = os.path.join(self.data_dir, filename)
            if os.path.isfile(file_path) and filename.lower().endswith((".pdf", ".pptx", ".txt")):
                print(f"Extracting text from: {filename}")
                final_chunks.extend(self.process_file(file_path))
        
        print(f"Successfully generated {len(final_chunks)} searchable chunks.")
        return final_chunks
