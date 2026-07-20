import os
import json
import re
import shutil
import textwrap
import time
from datetime import datetime
from pathlib import Path
from openpyxl import Workbook

from openpyxl.styles import Font, Alignment

import pandas as pd
from fastapi import FastAPI, File, HTTPException, UploadFile
from fastapi.responses import FileResponse
from fastapi.staticfiles import StaticFiles
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
from pydantic import BaseModel
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas

from agent_graph import app as graph_app
from agent_graph import db_manager
from document_processor import DocumentProcessor

BASE_DIR = Path(__file__).resolve().parent
DATA_DIR = BASE_DIR / "data"
OUTPUT_DIR = BASE_DIR / "outputs"
STATIC_DIR = BASE_DIR / "static"
SUPPORTED_DOCUMENT_EXTENSIONS = {".pdf", ".pptx", ".txt"}

for directory in [DATA_DIR, OUTPUT_DIR, STATIC_DIR]:
    directory.mkdir(exist_ok=True)

app = FastAPI(
    title="Agentic Document Intelligence System",
    version="2.0",
    description="RAG-based document assistant with local retrieval, web fallback, citations, and report exports.",
)


class QueryRequest(BaseModel):
    query: str
    output_format: str = "json"
    model_provider: str = "gemini"
    model_name: str | None = None
    temperature: float = 0.2
    max_tokens: int = 1400
    top_p: float = 0.9
    custom_answer: str | None = None
    presentation_mode: bool = False
    ppt_template: str = "strategy_brief"


class RefineRequest(BaseModel):
    query: str
    current_answer: str
    refinement_prompt: str
    model_provider: str = "gemini"
    model_name: str | None = None
    temperature: float = 0.2
    max_tokens: int = 1400
    top_p: float = 0.9


def _normalized_model_settings(request: QueryRequest):
    provider = request.model_provider.lower()
    if provider not in {"gemini", "openai", "groq"}:
        provider = "gemini"

    if provider == "openai":
        default_model = "gpt-4o-mini"
    elif provider == "groq":
        default_model = os.getenv("GROQ_MODEL", "llama-3.1-8b-instant")
    else:
        default_model = os.getenv("GEMINI_MODEL", "gemini-2.5-flash")
    max_token_limit = 1800 if provider == "groq" else 8000
    return {
        "model_provider": provider,
        "model_name": request.model_name or default_model,
        "temperature": max(0.0, min(2.0, request.temperature)),
        "max_tokens": max(256, min(max_token_limit, request.max_tokens)),
        "top_p": max(0.01, min(1.0, request.top_p)),
    }


def _timing_headers(timings):
    compact = {key: round(value, 4) for key, value in timings.items()}
    return {"X-Agent-Timings": json.dumps(compact, separators=(",", ":"))}


def _normalized_ppt_template(template):
    template = (template or "strategy_brief").strip().lower()
    return template if template in {"strategy_brief", "financial_review", "consulting_memo"} else "strategy_brief"


def _document_files():
    return sorted(
        file for file in DATA_DIR.iterdir()
        if file.is_file() and file.suffix.lower() in SUPPORTED_DOCUMENT_EXTENSIONS
    )


def _clear_uploaded_documents():
    deleted_files = []
    for file in _document_files():
        file.unlink()
        deleted_files.append(file.name)
    return deleted_files


def _safe_filename(name):
    keep = [char if char.isalnum() or char in "._- " else "_" for char in name]
    return "".join(keep).strip() or "uploaded_document"


def _report_name(query, extension):
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    slug = "".join(char.lower() if char.isalnum() else "_" for char in query[:40]).strip("_")
    return OUTPUT_DIR / f"agentic_report_{slug or 'query'}_{timestamp}.{extension}"


def _write_pdf_report(file_path, query, answer, sources):
    c = canvas.Canvas(str(file_path), pagesize=letter)
    width, height = letter
    y = height - 54

    c.setFont("Helvetica-Bold", 16)
    c.drawString(50, y, "Agentic Document Intelligence Report")
    y -= 28

    c.setFont("Helvetica", 10)
    c.drawString(50, y, f"Generated: {datetime.now().strftime('%d %b %Y, %I:%M %p')}")
    y -= 28

    c.setFont("Helvetica-Bold", 12)
    c.drawString(50, y, "Question")
    y -= 18
    c.setFont("Helvetica", 10)
    for line in textwrap.wrap(query, width=92):
        c.drawString(50, y, line)
        y -= 14

    y -= 14
    c.setFont("Helvetica-Bold", 12)
    c.drawString(50, y, "Answer")
    y -= 18
    c.setFont("Helvetica", 10)
    for paragraph in answer.splitlines() or [answer]:
        wrapped = textwrap.wrap(paragraph, width=92) or [""]
        for line in wrapped:
            if y < 70:
                c.showPage()
                y = height - 54
                c.setFont("Helvetica", 10)
            c.drawString(50, y, line)
            y -= 14

    if sources:
        y -= 12
        c.setFont("Helvetica-Bold", 12)
        c.drawString(50, y, "Sources")
        y -= 18
        c.setFont("Helvetica", 10)
        for source in sources:
            label = source.get("label") or source.get("source", "Unknown")
            for line in textwrap.wrap(f"- {label}", width=92):
                if y < 70:
                    c.showPage()
                    y = height - 54
                    c.setFont("Helvetica", 10)
                c.drawString(50, y, line)
                y -= 14

    c.save()





def _write_xlsx_report(file_path, query, answer, sources):
    wb = Workbook()
    ws = wb.active
    ws.title = "Agent Report"

    headers = [
        "User Query",
        "AI Response",
        "Sources",
        "Generated At"
    ]

    for col, header in enumerate(headers, start=1):
        cell = ws.cell(row=1, column=col)
        cell.value = header
        cell.font = Font(bold=True)

    source_text = "\n".join(
        source.get("label") or source.get("source", "")
        for source in sources
    )

    ws["A2"] = query
    ws["B2"] = answer
    ws["C2"] = source_text
    ws["D2"] = datetime.now().strftime("%d-%m-%Y %H:%M")

    # Wrap text
    for row in ws.iter_rows():
        for cell in row:
            cell.alignment = Alignment(
                wrap_text=True,
                vertical="top"
            )

    # Column widths
    ws.column_dimensions["A"].width = 35
    ws.column_dimensions["B"].width = 90
    ws.column_dimensions["C"].width = 45
    ws.column_dimensions["D"].width = 22

    # Increase row height
    ws.row_dimensions[2].height = 250

    wb.save(file_path)


def _answer_without_references(answer):
    content_lines = []
    for line in answer.splitlines():
        normalized = line.strip().lower()
        if normalized in {"sources:", "source:", "references:", "reference sources:", "reference links:"}:
            break
        content_lines.append(line)
    return "\n".join(content_lines).strip() or answer.strip()


def _clean_ppt_text(text):
    text = re.sub(r"\*\*\*(.*?)\*\*\*", r"\1", text)
    text = re.sub(r"\*\*(.*?)\*\*", r"\1", text)
    text = re.sub(r"\*(.*?)\*", r"\1", text)
    text = re.sub(r"`([^`]*)`", r"\1", text)
    text = re.sub(r"^#+\s*", "", text, flags=re.MULTILINE)
    text = text.replace("***", "")
    text = re.sub(r"[ \t]+", " ", text)
    return text.strip()


def _remove_low_value_lines(text):
    blocked_starts = (
        "the collected evidence indicates",
        "collected evidence indicates",
        "the provided evidence discusses",
        "the provided evidence does not",
        "the evidence discusses",
        "retrieved evidence discusses",
        "there is no specific document evidence",
        "the following answer is generated",
        "therefore, i will generate",
        "however, if you are asking",
        "if you are asking about",
        "generating such a portfolio would",
    )
    useful_lines = []
    for line in text.splitlines():
        stripped = line.strip()
        if not stripped:
            useful_lines.append("")
            continue
        if stripped.lower().startswith(blocked_starts):
            continue
        useful_lines.append(stripped)
    return "\n".join(useful_lines).strip()


def _presentation_text(answer):
    return _remove_low_value_lines(_clean_ppt_text(_answer_without_references(answer)))


def _split_sentences(text):
    normalized = re.sub(r"\s+", " ", text)
    sentences = re.split(r"(?<=[.!?])\s+(?=[A-Z0-9])", normalized)
    return [sentence.strip() for sentence in sentences if len(sentence.strip()) > 30]


def _slide_safe_phrase(text):
    clean = _clean_ppt_text(text)
    clean = re.sub(r"\s+", " ", clean).strip(" -")
    replacements = [
        (r"\bsmall to medium-sized businesses\b", "SMBs"),
        (r"\bsmall-to-medium-sized businesses\b", "SMBs"),
        (r"\bsmall and medium-sized businesses\b", "SMBs"),
        (r"\bThe company differentiates itself through its\b", "Differentiation comes from"),
        (r"\bDigitalOcean's pricing strategy is based on\b", "DigitalOcean pricing uses"),
        (r"\bThe company's target customers are those seeking\b", "Target customers seek"),
        (r"\boffering a unique value proposition of\b", "with a value proposition built on"),
        (r"\bwhich includes\b", "including"),
    ]
    for pattern, replacement in replacements:
        clean = re.sub(pattern, replacement, clean, flags=re.I)
    return clean.strip()


def _split_long_claim(text, limit=150):
    clean = _slide_safe_phrase(text).rstrip(".")
    if not clean:
        return []
    if len(clean) <= limit:
        return [clean + "."]

    split_patterns = [
        (r",\s+offering\s+", "Offers "),
        (r",\s+and\s+its\s+", "Its "),
        (r",\s+and\s+their\s+", "Their "),
        (r",\s+and\s+the\s+", "The "),
        (r",\s+while\s+", "While "),
        (r";\s+", ""),
        (r"\s+-\s+", ""),
    ]
    for pattern, right_prefix in split_patterns:
        match = re.search(pattern, clean, flags=re.I)
        if not match:
            continue
        left = clean[:match.start()].strip(" ,;:-")
        right = clean[match.end():].strip(" ,;:-")
        if not left or len(left) < 25 or not right or len(right) < 18:
            continue
        if right_prefix:
            right = right_prefix + right[0].lower() + right[1:] if right else right
        parts = _split_long_claim(left, limit=limit) + _split_long_claim(right, limit=limit)
        if len(parts) >= 2 and all(len(part) <= limit + 8 for part in parts):
            return parts

    comma_parts = [part.strip(" ,") for part in re.split(r",\s+", clean) if part.strip(" ,")]
    if len(comma_parts) > 2:
        chunks = []
        current = comma_parts[0]
        for part in comma_parts[1:]:
            candidate = f"{current}, {part}"
            if len(candidate) <= limit:
                current = candidate
            else:
                chunks.append(current)
                current = part
        chunks.append(current)
        if len(chunks) > 1 and all(len(chunk) >= 18 for chunk in chunks):
            return [chunk.rstrip(".") + "." for chunk in chunks[:3]]

    return [_shorten_text(clean, limit).rstrip(".") + "."]


def _dedupe_items(items, limit=None):
    seen = set()
    deduped = []
    for item in items:
        cleaned = _clean_ppt_text(item).strip(" -")
        cleaned = re.sub(r"\s+", " ", cleaned)
        if not cleaned:
            continue
        key = re.sub(r"[^a-z0-9]+", "", cleaned.lower())[:90]
        if key in seen:
            continue
        seen.add(key)
        deduped.append(cleaned)
        if limit and len(deduped) >= limit:
            break
    return deduped


def _extract_numbered_items(text):
    items = []
    pattern = re.compile(r"(?:^|\n)\s*(?:\d+[\.)]|[-•])\s+(.+?)(?=\n\s*(?:\d+[\.)]|[-•])\s+|\Z)", re.S)
    for match in pattern.finditer(text):
        item = " ".join(line.strip() for line in match.group(1).splitlines() if line.strip())
        item = re.sub(r"\s+", " ", item)
        if len(item) > 18:
            items.append(item)
    return _dedupe_items(items)


def _extract_named_sections(text):
    header_aliases = {
        "Market Portfolio": ["Market Portfolio", "DigitalOcean Market Portfolio"],
        "Company Overview": ["Company Overview"],
        "Financial Snapshot": ["Financial Snapshot", "Financial Overview", "Financial Performance", "Revenue And Profitability"],
        "Product Portfolio": ["Product Portfolio"],
        "Target Market Segments": ["Target Market Segments", "Target Segments", "Market Segments And Geographic Reach"],
        "Geographic Presence": ["Geographic Presence", "Geography"],
        "Competitive Landscape": ["Competitive Landscape", "Competition", "Competitive Landscape And Positioning"],
        "SWOT Analysis": ["SWOT Analysis", "SWOT"],
        "Market Positioning": ["Market Positioning"],
        "Customer Value Proposition": ["Customer Value Proposition", "Value Proposition"],
        "Revenue Streams": ["Revenue Streams", "Revenue And Go-To-Market"],
        "Marketing Strategy 4Ps": ["Marketing Strategy 4Ps", "Marketing Strategy (4Ps)", "Marketing Strategy"],
        "Growth Opportunities": ["Growth Opportunities", "Growth Opportunities And Strategic Priorities"],
        "Strategic Position Summary": ["Strategic Position Summary"],
        "Conclusion": ["Conclusion"],
        "Executive Summary": ["Executive Summary", "Summary"],
        "Market Position": ["Market Position"],
        "Portfolio Overview": ["Portfolio Overview", "Portfolio Snapshot"],
        "Comparative Analysis": ["Comparative Analysis", "Competitive Analysis"],
        "Strategic Recommendations": ["Strategic Recommendations", "Strategic Implications", "Strategic Moves"],
        "Suggested Table": ["Suggested Table", "Portfolio Analysis Table", "Segment Table"],
        "Risks And Mitigations": ["Risks And Mitigations", "Risks & Mitigations", "Risks"],
        "Success Metrics": ["Success Metrics", "Metrics"],
    }
    alias_to_header = {
        alias.lower(): header.lower()
        for header, aliases in header_aliases.items()
        for alias in aliases
    }
    aliases = sorted(alias_to_header, key=len, reverse=True)
    pattern = r"(?im)^\s*(?:\d+[\.)]\s*)?(?:" + "|".join(re.escape(alias) for alias in aliases) + r")\s*:?\s*$"
    matches = list(re.finditer(pattern, text))
    sections = {}

    for index, match in enumerate(matches):
        start = match.end()
        end = matches[index + 1].start() if index + 1 < len(matches) else len(text)
        matched_header = re.sub(r"^\d+[\.)]\s*", "", match.group(0).strip().rstrip(":")).lower()
        header = alias_to_header[matched_header]
        sections[header] = text[start:end].strip()

    return sections


def _section_bullets(section_text, limit=6):
    bullets = []
    for line in section_text.splitlines():
        stripped = line.strip()
        if not stripped:
            continue
        stripped = re.sub(r"^[-*•]\s*", "", stripped)
        stripped = re.sub(r"^\d+[\.)]\s*", "", stripped)
        stripped = re.sub(r"^[A-Z][A-Za-z /&-]{2,35}:\s*", "", stripped)
        stripped = re.sub(r"\s+", " ", stripped)
        if not stripped:
            continue
        if len(stripped) > 180 and not re.match(r"^[^:\n]{2,48}:", stripped):
            for sentence in _split_sentences(stripped):
                bullets.extend(_split_long_claim(sentence))
        else:
            bullets.extend(_split_long_claim(stripped))
    return _dedupe_items(bullets, limit=limit)


def _clean_analysis_item(item):
    item = _clean_ppt_text(item)
    item = re.sub(r"^[A-Z][A-Za-z /&-]{2,35}:\s*", "", item)
    item = re.sub(r"\b(Executive Summary|Market Position|Portfolio Overview|Comparative Analysis|Strategic Recommendations|Suggested Table|Risks And Mitigations|Success Metrics)\b", "", item)
    item = re.sub(r"\s+", " ", item).strip(" -:|")
    return item


def _parse_suggested_table(section_text):
    rows = []
    for line in section_text.splitlines():
        stripped = line.strip()
        if not stripped or "|" not in stripped:
            continue
        parts = [part.strip() for part in stripped.strip("|").split("|")]
        if all(re.fullmatch(r":?-{3,}:?", part.replace(" ", "")) for part in parts if part):
            continue
        if len(parts) >= 4:
            rows.append(parts[:4])

    if not rows:
        return {
            "headers": ["Segment", "Offering", "Customer Value", "Strategic Importance"],
            "rows": [],
        }

    header_row = rows[0]
    normalized_header = [cell.lower() for cell in header_row]
    expected = ["segment", "offering", "customer value", "strategic importance"]
    if normalized_header == expected:
        data_rows = rows[1:]
    else:
        data_rows = rows
        header_row = ["Segment", "Offering", "Customer Value", "Strategic Importance"]

    cleaned_rows = []
    for row in data_rows[:6]:
        cleaned_row = [_clean_ppt_text(cell)[:80] for cell in row]
        if any(cell for cell in cleaned_row):
            cleaned_rows.append(cleaned_row)

    return {"headers": header_row[:4], "rows": cleaned_rows}


def _parse_first_markdown_table(section_text):
    rows = []
    for line in section_text.splitlines():
        stripped = line.strip()
        if not stripped or "|" not in stripped:
            if rows:
                break
            continue
        parts = [part.strip() for part in stripped.strip("|").split("|")]
        if all(re.fullmatch(r":?-{3,}:?", part.replace(" ", "")) for part in parts if part):
            continue
        if len(parts) >= 2:
            rows.append(parts)

    if not rows:
        return None

    headers = [_clean_ppt_text(cell)[:42] for cell in rows[0]]
    financial_headers = {"metric", "latest value", "period / source", "business meaning"}
    max_rows = 5 if financial_headers.issubset({header.lower() for header in headers}) else 6
    data_rows = []
    for row in rows[1:1 + max_rows]:
        cleaned = [_clean_ppt_text(cell) for cell in row]
        if any(cleaned):
            data_rows.append(cleaned)

    return {"headers": headers, "rows": data_rows}


def _section_without_tables(section_text):
    lines = []
    in_table = False
    for line in section_text.splitlines():
        stripped = line.strip()
        if "|" in stripped:
            in_table = True
            continue
        if in_table and not stripped:
            in_table = False
            continue
        if not in_table:
            lines.append(line)
    return "\n".join(lines).strip()


def _short_phrase(text, fallback, limit=76):
    clean = _clean_ppt_text(text)
    clean = re.sub(r"^[^:]{1,35}:\s*", "", clean)
    return (clean or fallback)[:limit]


def _fallback_analysis_table(market_position, offerings, comparative, strategy):
    headers = ["Segment", "Offering", "Customer Value", "Strategic Importance"]
    segment_defaults = [
        "Developers / Builders",
        "Startups",
        "SMB Teams",
        "Platform Teams",
        "Growth Workloads",
        "Strategic Expansion",
    ]
    offering_defaults = [
        "Core cloud infrastructure",
        "Managed application services",
        "Developer workflow tooling",
        "Data and storage services",
        "Scalable deployment options",
        "Portfolio optimization",
    ]

    rows = []
    for index in range(6):
        offering_source = offerings[index % len(offerings)] if offerings else offering_defaults[index]
        value_source = market_position[index % len(market_position)] if market_position else offering_source
        importance_source = strategy[index % len(strategy)] if strategy else comparative[index % len(comparative)] if comparative else offering_source

        rows.append([
            segment_defaults[index],
            _short_phrase(offering_source, offering_defaults[index], limit=54),
            _short_phrase(value_source, "Clearer value for the target customer segment", limit=72),
            _short_phrase(importance_source, "Improves positioning, adoption, and portfolio growth", limit=78),
        ])

    return {"headers": headers, "rows": rows}


def _complete_analysis_table(table, market_position, offerings, comparative, strategy):
    fallback = _fallback_analysis_table(market_position, offerings, comparative, strategy)
    rows = list(table.get("rows") or [])
    seen = {"|".join(cell.lower() for cell in row) for row in rows}

    for row in fallback["rows"]:
        key = "|".join(cell.lower() for cell in row)
        if key in seen:
            continue
        rows.append(row)
        seen.add(key)
        if len(rows) >= 6:
            break

    return {
        "headers": table.get("headers") or fallback["headers"],
        "rows": rows[:6],
    }


def _item_key(item):
    return re.sub(r"[^a-z0-9]+", "", _clean_ppt_text(item).lower())[:100]


def _take_unused(items, used, limit):
    selected = []
    for item in items:
        key = _item_key(item)
        if not key or key in used:
            continue
        used.add(key)
        selected.append(item)
        if len(selected) >= limit:
            break
    return selected


def _fallback_section_buckets(clean_text):
    all_points = [
        _clean_analysis_item(point)
        for point in _extract_numbered_items(clean_text) + _split_sentences(clean_text)
    ]
    all_points = _dedupe_items(all_points)
    if not all_points:
        all_points = [
            "The generated portfolio summarizes the market opportunity, target customers, offering mix, and strategic considerations.",
            "The portfolio should be reviewed against current market evidence before final business decisions are made.",
        ]

    used = set()
    summary = _take_unused(all_points, used, 4)
    market_position = _take_unused(all_points, used, 5)
    offerings = _take_unused(all_points, used, 6)
    comparative = _take_unused(all_points, used, 5)
    strategy = _take_unused(all_points, used, 5)
    risks = _take_unused(all_points, used, 4)
    metrics = _take_unused(all_points, used, 5)

    # If the draft is too short, fill the remaining slides with purposeful defaults, not reused text.
    market_position = market_position or [
        "Target customers, differentiation, and competitive position need to be evaluated together to avoid a purely descriptive portfolio."
    ]
    offerings = offerings or [
        "Core offerings should be grouped by customer need, commercial value, and strategic importance."
    ]
    comparative = comparative or [
        "The portfolio should compare strengths, gaps, and alternatives rather than listing capabilities in isolation."
    ]
    strategy = strategy or [
        "Recommended actions should connect the portfolio to customer acquisition, retention, and expansion outcomes."
    ]
    risks = risks or [
        "Main risks include competitive pressure, weak differentiation, and execution gaps; mitigation should focus on sharper positioning."
    ]
    metrics = metrics or [
        "Track adoption, conversion, retention, usage depth, and expansion to judge portfolio performance."
    ]

    return {
        "summary": summary[:4],
        "market_position": market_position[:5],
        "offerings": offerings[:6],
        "comparative": comparative[:5],
        "strategy": strategy[:5],
        "risks": risks[:4],
        "metrics": metrics[:5],
    }


def _default_summary_points(query):
    topic = query.strip().rstrip(".") or "the requested portfolio"
    topic = re.sub(r"(?i)^\s*(generate|create|make|prepare|build)\s+(the\s+)?", "", topic).strip()
    topic = topic or "the requested portfolio"
    return [
        f"The {topic} should be evaluated through customer fit, differentiated positioning, portfolio strength, and execution risk.",
        "The main opportunity is to connect offerings to clear customer segments instead of presenting capabilities as an isolated list.",
        "Competitive analysis should explain where the portfolio is defensible, where it is weaker, and what alternatives customers may choose.",
        "The recommended direction should translate market position into practical product, go-to-market, and measurement priorities.",
    ]


def _make_table_slide(title, subtitle, section_text, kicker):
    table = _parse_first_markdown_table(section_text)
    if not table:
        return {
            "kind": "bullets",
            "title": title,
            "subtitle": subtitle,
            "bullets": _section_bullets(section_text, limit=6),
            "kicker": kicker,
        }
    if len(table["rows"]) <= 2:
        card_items = []
        for row in table["rows"]:
            for index, header in enumerate(table["headers"]):
                value = row[index] if index < len(row) else ""
                if value:
                    card_items.append(f"{header}: {value}")
        return {
            "kind": "cards",
            "title": title,
            "subtitle": subtitle,
            "items": card_items,
            "kicker": kicker,
        }
    return {
        "kind": "table",
        "title": title,
        "subtitle": subtitle,
        "table": table,
        "kicker": kicker,
    }


def _swot_groups(section_text):
    groups = {}
    current = None
    for line in section_text.splitlines():
        stripped = line.strip()
        if not stripped:
            continue
        normalized = stripped.strip(":").lower()
        if normalized in {"strengths", "weaknesses", "opportunities", "threats"}:
            current = normalized.title()
            groups[current] = []
            continue
        if current:
            cleaned = re.sub(r"^[-*•]\s*", "", stripped)
            cleaned = re.sub(r"^\d+[\.)]\s*", "", cleaned)
            cleaned = _clean_analysis_item(cleaned)
            if cleaned:
                groups[current].append(cleaned)
    return groups


def _swot_bullets(section_text):
    groups = _swot_groups(section_text)
    bullets = []
    for label in ["Strengths", "Weaknesses", "Opportunities", "Threats"]:
        items = groups.get(label, [])[:3]
        if items:
            bullets.append(f"{label}: " + "; ".join(items))
    return bullets or _section_bullets(section_text, limit=6)


def _table_rows_as_bullets(section_text, limit=5):
    table = _parse_first_markdown_table(section_text)
    if not table:
        return _section_bullets(section_text, limit=limit)

    headers = table["headers"]
    bullets = []
    for row in table["rows"][:limit]:
        cells = [cell for cell in row if cell]
        if not cells:
            continue
        if len(cells) >= 3:
            bullets.append(f"{cells[0]}: {cells[1]} - {cells[2]}")
        elif len(cells) == 2:
            bullets.append(f"{cells[0]}: {cells[1]}")
        else:
            bullets.append(cells[0])

    if not bullets and headers:
        bullets = _section_bullets(section_text, limit=limit)
    return bullets[:limit]


def _combine_bullets(*sections, limit=6):
    combined = []
    for section in sections:
        if not section:
            continue
        combined.extend(_table_rows_as_bullets(section, limit=limit))
        combined.extend(_section_bullets(_section_without_tables(section), limit=limit))
    return _dedupe_items(combined, limit=limit)


def _build_market_portfolio_slides(named_sections):
    market_keys = {
        "company overview",
        "financial snapshot",
        "product portfolio",
        "target market segments",
        "competitive landscape",
        "swot analysis",
        "revenue streams",
        "marketing strategy 4ps",
        "strategic position summary",
    }
    if not any(key in named_sections for key in market_keys):
        return []

    slides = []
    if "company overview" in named_sections:
        slides.append(_make_table_slide(
            "Company Snapshot Establishes The Business Context",
            "Audience, category, and core value proposition",
            named_sections["company overview"],
            "overview",
        ))
    if "financial snapshot" in named_sections:
        slides.append(_make_table_slide(
            "Financial Snapshot Quantifies Scale And Momentum",
            "Revenue, profitability, growth, and operating signals",
            named_sections["financial snapshot"],
            "metrics",
        ))
    if "product portfolio" in named_sections:
        slides.append(_make_table_slide(
            "Product Portfolio Maps Capabilities To Business Roles",
            "Core categories and why each matters",
            named_sections["product portfolio"],
            "portfolio",
        ))

    market_reach = _combine_bullets(
        named_sections.get("target market segments", ""),
        named_sections.get("geographic presence", ""),
        limit=6,
    )
    if market_reach:
        slides.append({
            "kind": "bullets",
            "title": "Target Market And Reach Define The Opportunity",
            "subtitle": "Customer segments, needs, and geographic logic",
            "bullets": market_reach,
            "kicker": "market",
        })

    if "competitive landscape" in named_sections:
        slides.append(_make_table_slide(
            "Competition And Positioning Clarify Differentiation",
            "Competitor strengths, relative gaps, and implications",
            named_sections["competitive landscape"],
            "competition",
        ))
    if "swot analysis" in named_sections:
        slides.append({
            "kind": "swot",
            "title": "SWOT Analysis Frames Internal And External Factors",
            "subtitle": "Strengths, weaknesses, opportunities, and threats",
            "swot": _swot_groups(named_sections["swot analysis"]),
            "kicker": "swot",
        })

    value_points = _combine_bullets(
        named_sections.get("market positioning", ""),
        named_sections.get("customer value proposition", ""),
        limit=6,
    )
    if value_points:
        slides.append({
            "kind": "bullets",
            "title": "Positioning And Value Proposition Explain The Niche",
            "subtitle": "Practical customer benefits and buying logic",
            "bullets": value_points,
            "kicker": "value",
        })

    revenue_gtm = _combine_bullets(
        named_sections.get("revenue streams", ""),
        named_sections.get("marketing strategy 4ps", ""),
        limit=6,
    )
    if revenue_gtm:
        slides.append({
            "kind": "bullets",
            "title": "Revenue And Go-To-Market Connect Strategy To Execution",
            "subtitle": "Monetization, pricing, distribution, and promotion logic",
            "bullets": revenue_gtm,
            "kicker": "gtm",
        })

    growth_strategy = _combine_bullets(
        named_sections.get("growth opportunities", ""),
        named_sections.get("strategic position summary", ""),
        named_sections.get("conclusion", ""),
        limit=6,
    )
    if growth_strategy:
        slides.append({
            "kind": "bullets",
            "title": "Growth Priorities Summarize The Strategic Path",
            "subtitle": "Expansion moves, ratings, and final recommendation",
            "bullets": growth_strategy,
            "kicker": "growth",
        })

    if "risks and mitigations" in named_sections:
        slides.append({
            "kind": "risk_matrix",
            "title": "Risks Need Practical Mitigation Paths",
            "subtitle": "Key risks and the actions that reduce them",
            "bullets": _section_bullets(named_sections["risks and mitigations"], limit=4),
            "kicker": "risk",
        })

    if "success metrics" in named_sections:
        slides.append({
            "kind": "kpi",
            "title": "Success Metrics Show Whether The Strategy Works",
            "subtitle": "Indicators to track adoption, retention, and growth",
            "bullets": _section_bullets(named_sections["success metrics"], limit=5),
            "kicker": "metrics",
        })

    return slides[:10]


def _build_ppt_sections(query, answer):
    clean_text = _presentation_text(answer)
    named_sections = _extract_named_sections(clean_text)
    if named_sections:
        extra_slides = _build_market_portfolio_slides(named_sections)
        title = query.strip().rstrip(".")
        if title:
            title = title[0].upper() + title[1:]
        else:
            title = "Generated Market Portfolio"

        summary = _section_bullets(named_sections.get("executive summary", ""), limit=4)
        market_position = _section_bullets(named_sections.get("market position", ""), limit=5)
        offerings = _section_bullets(named_sections.get("portfolio overview", ""), limit=6)
        comparative = _section_bullets(named_sections.get("comparative analysis", ""), limit=5)
        strategy = _section_bullets(named_sections.get("strategic recommendations", ""), limit=5)
        table = _parse_suggested_table(named_sections.get("suggested table", ""))
        risks = _section_bullets(named_sections.get("risks and mitigations", ""), limit=4)
        metrics = _section_bullets(named_sections.get("success metrics", ""), limit=5)
        fallback_text = "\n\n".join(section for section in named_sections.values() if section.strip())
        fallback = _fallback_section_buckets(fallback_text or clean_text)

        used = set()
        summary = summary or _default_summary_points(query)
        market_position = market_position or _take_unused(fallback["market_position"], used, 5)
        offerings = offerings or _take_unused(fallback["offerings"], used, 6)
        comparative = comparative or _take_unused(fallback["comparative"], used, 5)
        strategy = strategy or _take_unused(fallback["strategy"], used, 5)
        risks = risks or _take_unused(fallback["risks"], used, 4)
        metrics = metrics or _take_unused(fallback["metrics"], used, 5)
        table = _complete_analysis_table(table, market_position, offerings, comparative, strategy)

        return {
            "title": title,
            "scope": [
                "The portfolio frames the market opportunity, target customers, offering mix, and strategic actions.",
                "The analysis separates core business content from references so the main story stays presentation-ready.",
                "The deck is structured for mentor review with summary, positioning, comparison, risks, and metrics.",
            ],
            "summary": summary,
            "market_position": market_position,
            "offerings": offerings,
            "comparative": comparative,
            "strategy": strategy,
            "table": table,
            "risks": risks,
            "metrics": metrics,
            "extra_slides": extra_slides,
        }

    buckets = _fallback_section_buckets(clean_text)

    title = query.strip().rstrip(".")
    if title:
        title = title[0].upper() + title[1:]
    else:
        title = "Generated Market Portfolio"

    return {
        "title": title,
        "scope": [
            "The portfolio frames the market opportunity, target customers, offering mix, and strategic actions.",
            "The analysis separates core business content from references so the main story stays presentation-ready.",
            "The deck is structured for mentor review with summary, positioning, comparison, risks, and metrics.",
        ],
        "summary": buckets["summary"],
        "market_position": buckets["market_position"],
        "offerings": buckets["offerings"],
        "comparative": buckets["comparative"],
        "strategy": buckets["strategy"],
        "table": _fallback_analysis_table(
            buckets["market_position"],
            buckets["offerings"],
            buckets["comparative"],
            buckets["strategy"],
        ),
        "risks": buckets["risks"],
        "metrics": buckets["metrics"],
        "extra_slides": [],
    }


def _source_label(source):
    label = source.get("label") or source.get("source") or "Unknown source"
    url = source.get("url")
    if url:
        return f"{label} - {url}"
    return label


def _set_run_style(run, size=18, bold=False, color=(23, 32, 44)):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*color)
    run.font.name = "Aptos"


PPT_NAVY = (16, 32, 52)
PPT_BLUE = (25, 103, 210)
PPT_TEAL = (41, 171, 135)
PPT_INK = (31, 42, 58)
PPT_MUTED = (96, 112, 132)
PPT_LINE = (218, 226, 236)
PPT_BG = (247, 250, 253)
PPT_SOFT_BLUE = (229, 240, 255)
PPT_SOFT_TEAL = (229, 247, 242)
PPT_WARN = (190, 82, 65)
PPT_GREEN = (36, 145, 105)
PPT_INDIGO = (79, 70, 229)
PPT_EMERALD = (5, 150, 105)
PPT_SLATE = (30, 41, 59)
PPT_LIGHT_INDIGO = (238, 242, 255)
PPT_LIGHT_EMERALD = (236, 253, 245)

PPT_THEMES = {
    "strategy_brief": {
        "name": "Classic Consulting Layout",
        "accent": PPT_BLUE,
        "secondary": PPT_TEAL,
        "navy": PPT_NAVY,
        "bg": PPT_BG,
        "muted": PPT_MUTED,
        "soft": PPT_SOFT_BLUE,
        "style": "classic",
    },
    "financial_review": {
        "name": "Split-Column Layout",
        "accent": PPT_EMERALD,
        "secondary": (14, 116, 144),
        "navy": (15, 46, 43),
        "bg": (246, 252, 249),
        "muted": (81, 105, 99),
        "soft": PPT_LIGHT_EMERALD,
        "style": "sidebar",
    },
    "consulting_memo": {
        "name": "Card-Based Layout",
        "accent": PPT_INDIGO,
        "secondary": (14, 165, 233),
        "navy": PPT_SLATE,
        "bg": (248, 250, 255),
        "muted": (99, 102, 129),
        "soft": PPT_LIGHT_INDIGO,
        "style": "editorial",
    },
}

SECTION_ICONS = {
    "overview": "◼",
    "portfolio": "☁",
    "market": "👥",
    "segments": "👥",
    "geography": "◎",
    "competition": "↔",
    "comparison": "↔",
    "swot": "◆",
    "position": "◇",
    "value": "★",
    "gtm": "▶",
    "strategy": "▶",
    "growth": "↗",
    "risk": "!",
    "metrics": "▣",
    "table": "▦",
    "summary": "●",
    "sources": "i",
    "context": "●",
}


def _shorten_text(text, limit=145):
    clean = _clean_ppt_text(text)
    clean = re.sub(r"\s+", " ", clean).strip()
    if len(clean) <= limit:
        return clean

    sentences = re.split(r"(?<=[.!?])\s+", clean)
    if sentences and len(sentences[0]) <= limit:
        return sentences[0].rstrip(".")

    clauses = re.split(r"\s*;\s*", clean)
    if len(clauses) > 1:
        selected = []
        current = ""
        for clause in clauses:
            candidate = "; ".join([*selected, clause])
            if len(candidate) > limit:
                break
            selected.append(clause)
            current = candidate
        if current:
            return current.rstrip(" .;")

    dash_parts = re.split(r"\s+-\s+", clean)
    if len(dash_parts) > 1 and len(dash_parts[0]) <= limit:
        return dash_parts[0].rstrip(" .;-")

    cut = clean[:limit].rsplit(" ", 1)[0].rstrip(" ,.;:")
    return cut


def _compact_bullets(items, limit=6, char_limit=145):
    compact = []
    for item in items:
        shortened = _shorten_text(item, char_limit)
        if shortened:
            compact.append(shortened)
        if len(compact) >= limit:
            break
    return compact


def _shorten_table_text(text, limit=72):
    clean = _clean_ppt_text(text)
    clean = re.sub(r"\s+", " ", clean).strip()
    if len(clean) <= limit:
        return clean

    sentence_parts = re.split(r"(?<=[.!?])\s+", clean)
    if sentence_parts and 18 <= len(sentence_parts[0]) <= limit:
        return sentence_parts[0].rstrip(".")

    pre_replacements = [
        (r"\bexceeding initial guidance\b", "beating guidance"),
        (r"\bdemonstrating strong execution on growth strategy\b", "showing strong execution"),
        (r"\bindicating strong market momentum\b", "showing market momentum"),
        (r"\bindicating operational efficiency and scalability\b", "showing efficiency and scale"),
    ]
    for pattern, replacement in pre_replacements:
        clean = re.sub(pattern, replacement, clean, flags=re.I)
    if len(clean) <= limit:
        return clean

    clauses = re.split(
        r"\s*(?:;|, and|, but|, while|, exceeding|, indicating|, demonstrating|, driving)\s*",
        clean,
        flags=re.I,
    )
    clauses = [clause.strip(" ,.;:-") for clause in clauses if clause.strip(" ,.;:-")]
    if clauses:
        selected = []
        current = ""
        for clause in clauses:
            candidate = "; ".join([*selected, clause])
            if len(candidate) > limit:
                break
            selected.append(clause)
            current = candidate
        if current and len(current) >= 18:
            return current.rstrip(" ,.;:-")

    compact_replacements = [
        (r"\bAccelerating\b", "Fast"),
        (r"\bdemonstrating\b", "showing"),
        (r"\bindicating\b", "showing"),
        (r"\bcompared to\b", "vs."),
        (r"\boperational efficiency\b", "efficiency"),
        (r"\bmarket leadership\b", "leadership"),
        (r"\bcustomer acquisition\b", "acquisition"),
        (r"\bcustomer retention\b", "retention"),
    ]
    compact = clean
    for pattern, replacement in compact_replacements:
        compact = re.sub(pattern, replacement, compact, flags=re.I)
    if len(compact) <= limit:
        return compact

    cut = compact[:limit].rsplit(" ", 1)[0].rstrip(" ,.;:-")
    trailing_fragments = {
        "and", "or", "but", "while", "with", "without", "through", "from",
        "to", "of", "on", "for", "initial", "strong", "market", "growth",
    }
    words = cut.split()
    while words and words[-1].lower() in trailing_fragments:
        words.pop()
    return " ".join(words) if words else cut


def _split_label(text):
    match = re.match(r"^([^:\n]{2,48}):\s*(.+)$", text)
    if not match:
        return None, text
    return match.group(1), match.group(2)


def _ppt_theme(template):
    return PPT_THEMES.get(_normalized_ppt_template(template), PPT_THEMES["strategy_brief"])


def _add_label_detail_runs(paragraph, text, size, label_color=PPT_BLUE, detail_color=PPT_INK):
    label, detail = _split_label(text)
    if label:
        label_run = paragraph.add_run()
        label_run.text = f"{label}: "
        _set_run_style(label_run, size=size, bold=True, color=label_color)

        detail_run = paragraph.add_run()
        detail_run.text = detail
        _set_run_style(detail_run, size=size, color=detail_color)
    else:
        run = paragraph.add_run()
        run.text = text
        _set_run_style(run, size=size, color=detail_color)


def _apply_slide_background(slide, template="strategy_brief"):
    theme = _ppt_theme(template)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*theme["bg"])

    if theme["style"] == "sidebar":
        side = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.42), Inches(7.5))
        side.fill.solid()
        side.fill.fore_color.rgb = RGBColor(*theme["navy"])
        side.line.fill.background()

        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.42), Inches(0), Inches(0.08), Inches(7.5))
        accent.fill.solid()
        accent.fill.fore_color.rgb = RGBColor(*theme["accent"])
        accent.line.fill.background()
        return

    top_height = 0.16 if theme["style"] == "classic" else 0.52
    top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(top_height))
    top.fill.solid()
    top.fill.fore_color.rgb = RGBColor(*theme["navy"])
    top.line.fill.background()

    accent_y = top_height
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(accent_y), Inches(13.333), Inches(0.045))
    accent.fill.solid()
    accent.fill.fore_color.rgb = RGBColor(*theme["accent"])
    accent.line.fill.background()


def _text_size_for_count(count, base=16):
    if count <= 4:
        return base + 3
    if count == 5:
        return base + 2
    return base + 1


def _add_deck_label(slide, template="strategy_brief"):
    theme = _ppt_theme(template)
    label = slide.shapes.add_textbox(Inches(10.35), Inches(0.36), Inches(2.32), Inches(0.28))
    frame = label.text_frame
    frame.clear()
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.RIGHT
    run = paragraph.add_run()
    run.text = "Agentic Document Intelligence"
    color = (235, 242, 255) if theme["style"] == "editorial" else theme["muted"]
    _set_run_style(run, size=8.5, bold=True, color=color)


def _add_section_kicker(slide, text, template="strategy_brief"):
    theme = _ppt_theme(template)
    x = 0.78 if theme["style"] != "sidebar" else 0.72
    y = 0.34 if theme["style"] != "editorial" else 0.68
    chip = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(1.55), Inches(0.32))
    chip.fill.solid()
    chip.fill.fore_color.rgb = RGBColor(*theme["soft"])
    chip.line.color.rgb = RGBColor(*theme["accent"])

    box = slide.shapes.add_textbox(Inches(x + 0.1), Inches(y + 0.05), Inches(1.35), Inches(0.2))
    frame = box.text_frame
    frame.clear()
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = text.upper()
    _set_run_style(run, size=8.5, bold=True, color=theme["accent"])


def _add_title(slide, title, subtitle=None, kicker="analysis", template="strategy_brief"):
    theme = _ppt_theme(template)
    _add_section_kicker(slide, kicker, template)
    _add_deck_label(slide, template)

    title_x = 1.28
    icon_x = 0.72
    title_y = 0.66
    if theme["style"] == "sidebar":
        title_x = 1.52
        icon_x = 0.78
    elif theme["style"] == "editorial":
        title_y = 0.9

    icon_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(icon_x), Inches(title_y + 0.1), Inches(0.46), Inches(0.46))
    icon_box.fill.solid()
    icon_box.fill.fore_color.rgb = RGBColor(*theme["accent"])
    icon_box.line.fill.background()

    icon_text = slide.shapes.add_textbox(Inches(icon_x), Inches(title_y + 0.165), Inches(0.46), Inches(0.22))
    icon_frame = icon_text.text_frame
    icon_frame.clear()
    icon_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    icon_run = icon_frame.paragraphs[0].add_run()
    icon_run.text = SECTION_ICONS.get(kicker, "●")
    _set_run_style(icon_run, size=13, bold=True, color=(255, 255, 255))

    title_box = slide.shapes.add_textbox(Inches(title_x), Inches(title_y), Inches(10.25), Inches(1.16))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_frame.clear()
    title_paragraph = title_frame.paragraphs[0]
    title_run = title_paragraph.add_run()
    clean_title = _shorten_text(title, 64)
    title_run.text = clean_title
    if len(clean_title) > 50:
        title_size = 29
    elif len(clean_title) > 42:
        title_size = 32
    else:
        title_size = 36
    _set_run_style(title_run, size=title_size, bold=True, color=theme["navy"])

    if subtitle:
        subtitle_y = 2.0 if theme["style"] != "editorial" else 2.08
        subtitle_box = slide.shapes.add_textbox(Inches(title_x + 0.02), Inches(subtitle_y), Inches(10.95), Inches(0.4))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.clear()
        subtitle_run = subtitle_frame.paragraphs[0].add_run()
        subtitle_run.text = subtitle[:145]
        _set_run_style(subtitle_run, size=16, color=theme["muted"])

    divider_x = 0.74 if theme["style"] != "sidebar" else 0.82
    divider = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(divider_x), Inches(2.32), Inches(11.85), Inches(0.025))
    divider.fill.solid()
    divider.fill.fore_color.rgb = RGBColor(*theme["accent"] if theme["style"] != "classic" else PPT_LINE)
    divider.line.fill.background()


def _add_body(slide, text, top=1.62, height=4.95, size=15):
    body_box = slide.shapes.add_textbox(Inches(0.78), Inches(top), Inches(11.85), Inches(height))
    body_frame = body_box.text_frame
    body_frame.word_wrap = True
    body_frame.clear()

    paragraphs = [paragraph.strip() for paragraph in text.split("\n\n") if paragraph.strip()]
    for index, paragraph_text in enumerate(paragraphs):
        paragraph = body_frame.paragraphs[0] if index == 0 else body_frame.add_paragraph()
        paragraph.space_after = Pt(8)
        run = paragraph.add_run()
        run.text = paragraph_text[:650]
        _set_run_style(run, size=size, color=(39, 50, 68))


def _add_cover_body(slide, text, x=0.88, y=5.55, color=(225, 235, 247)):
    for index, line in enumerate(text.splitlines()):
        paragraph_text = line if index == 0 else f"Prepared {line}"
        box = slide.shapes.add_textbox(Inches(x), Inches(y + index * 0.32), Inches(5.9), Inches(0.3))
        frame = box.text_frame
        frame.clear()
        run = frame.paragraphs[0].add_run()
        run.text = paragraph_text
        _set_run_style(run, size=13 if index == 0 else 11.5, bold=index == 0, color=color)


def _add_cover_slide(slide, title, subtitle, body, template="strategy_brief"):
    theme = _ppt_theme(template)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*theme["navy"])

    if theme["style"] == "sidebar":
        fill.fore_color.rgb = RGBColor(245, 252, 249)
        side = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(3.45), Inches(7.5))
        side.fill.solid()
        side.fill.fore_color.rgb = RGBColor(*theme["navy"])
        side.line.fill.background()

        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.45), Inches(0), Inches(0.12), Inches(7.5))
        accent.fill.solid()
        accent.fill.fore_color.rgb = RGBColor(*theme["accent"])
        accent.line.fill.background()

        title_box = slide.shapes.add_textbox(Inches(4.05), Inches(1.28), Inches(7.9), Inches(2.05))
        frame = title_box.text_frame
        frame.word_wrap = True
        frame.clear()
        run = frame.paragraphs[0].add_run()
        run.text = title[:95]
        _set_run_style(run, size=42, bold=True, color=theme["navy"])

        subtitle_box = slide.shapes.add_textbox(Inches(4.08), Inches(3.5), Inches(6.8), Inches(0.48))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.clear()
        subtitle_run = subtitle_frame.paragraphs[0].add_run()
        subtitle_run.text = subtitle
        _set_run_style(subtitle_run, size=18, bold=True, color=theme["accent"])

        callout = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.08), Inches(4.35), Inches(6.65), Inches(1.05))
        callout.fill.solid()
        callout.fill.fore_color.rgb = RGBColor(255, 255, 255)
        callout.line.color.rgb = RGBColor(*PPT_LINE)
        callout_text = slide.shapes.add_textbox(Inches(4.35), Inches(4.58), Inches(6.05), Inches(0.46))
        callout_frame = callout_text.text_frame
        callout_frame.clear()
        callout_run = callout_frame.paragraphs[0].add_run()
        callout_run.text = "Same analysis, structured in split-column executive pages."
        _set_run_style(callout_run, size=16, bold=True, color=theme["navy"])
        _add_cover_body(slide, body, x=0.55, y=5.7, color=(220, 242, 235))
        return

    if theme["style"] == "editorial":
        fill.fore_color.rgb = RGBColor(248, 250, 255)
        header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.2))
        header.fill.solid()
        header.fill.fore_color.rgb = RGBColor(*theme["navy"])
        header.line.fill.background()

        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.82), Inches(1.52), Inches(1.1), Inches(0.1))
        accent.fill.solid()
        accent.fill.fore_color.rgb = RGBColor(*theme["accent"])
        accent.line.fill.background()

        title_box = slide.shapes.add_textbox(Inches(0.82), Inches(1.78), Inches(7.3), Inches(2.1))
        frame = title_box.text_frame
        frame.word_wrap = True
        frame.clear()
        run = frame.paragraphs[0].add_run()
        run.text = title[:95]
        _set_run_style(run, size=42, bold=True, color=theme["navy"])

        subtitle_box = slide.shapes.add_textbox(Inches(0.86), Inches(4.02), Inches(5.9), Inches(0.46))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.clear()
        subtitle_run = subtitle_frame.paragraphs[0].add_run()
        subtitle_run.text = subtitle
        _set_run_style(subtitle_run, size=18, bold=True, color=theme["accent"])

        for idx, label in enumerate(["Cards", "KPIs", "Tables"]):
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.25), Inches(1.75 + idx * 1.18), Inches(3.6), Inches(0.82))
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(255, 255, 255)
            card.line.color.rgb = RGBColor(*PPT_LINE)
            box = slide.shapes.add_textbox(Inches(8.48), Inches(1.98 + idx * 1.18), Inches(3.0), Inches(0.24))
            frame = box.text_frame
            frame.clear()
            run = frame.paragraphs[0].add_run()
            run.text = label
            _set_run_style(run, size=15, bold=True, color=theme["navy"])
        _add_cover_body(slide, body, x=0.88, y=5.52, color=theme["muted"])
        return

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.22), Inches(7.5))
    accent.fill.solid()
    accent.fill.fore_color.rgb = RGBColor(*theme["secondary"])
    accent.line.fill.background()

    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.85), Inches(0), Inches(5.48), Inches(7.5))
    band.fill.solid()
    band.fill.fore_color.rgb = RGBColor(21, 45, 73)
    band.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.82), Inches(1.0), Inches(6.8), Inches(2.2))
    frame = title_box.text_frame
    frame.word_wrap = True
    frame.clear()
    run = frame.paragraphs[0].add_run()
    run.text = title[:95]
    _set_run_style(run, size=43, bold=True, color=(255, 255, 255))

    subtitle_box = slide.shapes.add_textbox(Inches(0.88), Inches(3.28), Inches(6.2), Inches(0.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.clear()
    subtitle_run = subtitle_frame.paragraphs[0].add_run()
    subtitle_run.text = subtitle
    _set_run_style(subtitle_run, size=18, bold=True, color=(161, 220, 205))

    _add_cover_body(slide, body)

    callout = slide.shapes.add_textbox(Inches(8.25), Inches(1.35), Inches(4.35), Inches(3.7))
    callout_frame = callout.text_frame
    callout_frame.word_wrap = True
    callout_frame.clear()
    heading = callout_frame.paragraphs[0].add_run()
    heading.text = "Presentation-ready analysis"
    _set_run_style(heading, size=26, bold=True, color=(255, 255, 255))
    paragraph = callout_frame.add_paragraph()
    paragraph.space_before = Pt(16)
    detail = paragraph.add_run()
    detail.text = "Executive summary, market position, portfolio view, analysis table, risks, metrics, and separated references."
    _set_run_style(detail, size=18, color=(211, 224, 240))


def _add_bullets(slide, bullets, top=2.58, left=0.82, width=11.75, size=None):
    raw_bullets = [bullet for bullet in bullets if bullet]
    char_limit = 190 if len(raw_bullets) <= 4 else 155
    bullets = _compact_bullets(raw_bullets, limit=6, char_limit=char_limit)
    size = size or _text_size_for_count(len(bullets), 16)
    bullet_box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(4.1))
    bullet_frame = bullet_box.text_frame
    bullet_frame.word_wrap = True
    bullet_frame.margin_left = Inches(0.08)
    bullet_frame.margin_right = Inches(0.08)
    bullet_frame.clear()

    for index, bullet in enumerate(bullets):
        paragraph = bullet_frame.paragraphs[0] if index == 0 else bullet_frame.add_paragraph()
        paragraph.level = 0
        paragraph.space_after = Pt(13)
        _add_label_detail_runs(paragraph, bullet, size=size, label_color=PPT_BLUE)


def _add_bullets_two_column(slide, bullets, template="financial_review"):
    theme = _ppt_theme(template)
    raw_bullets = [bullet for bullet in bullets if bullet]
    bullets = _compact_bullets(raw_bullets, limit=6, char_limit=135)
    columns = [bullets[:3], bullets[3:6]]
    positions = [(0.82, 2.58), (6.82, 2.58)]
    for column_index, column_items in enumerate(columns):
        if not column_items:
            continue
        x, y = positions[column_index]
        rail = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.06), Inches(3.82))
        rail.fill.solid()
        rail.fill.fore_color.rgb = RGBColor(*(theme["accent"] if column_index == 0 else theme["secondary"]))
        rail.line.fill.background()

        box = slide.shapes.add_textbox(Inches(x + 0.22), Inches(y), Inches(5.35), Inches(3.9))
        frame = box.text_frame
        frame.word_wrap = True
        frame.margin_left = Inches(0.04)
        frame.margin_right = Inches(0.04)
        frame.clear()
        for index, bullet in enumerate(column_items):
            paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
            paragraph.space_after = Pt(15)
            _add_label_detail_runs(paragraph, bullet, size=15.5, label_color=theme["accent"] if column_index == 0 else theme["secondary"])


def _add_bullet_cards(slide, bullets, template="consulting_memo"):
    theme = _ppt_theme(template)
    items = _compact_bullets([bullet for bullet in bullets if bullet], limit=6, char_limit=110)
    positions = [
        (0.82, 2.55), (4.72, 2.55), (8.62, 2.55),
        (0.82, 4.32), (4.72, 4.32), (8.62, 4.32),
    ]
    card_width = 3.45
    card_height = 1.34
    for index, item in enumerate(items):
        x, y = positions[index]
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(card_width), Inches(card_height))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(255, 255, 255)
        card.line.color.rgb = RGBColor(*theme["soft"])

        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.09), Inches(card_height))
        accent.fill.solid()
        accent.fill.fore_color.rgb = RGBColor(*(theme["accent"] if index % 2 == 0 else theme["secondary"]))
        accent.line.fill.background()

        text_box = slide.shapes.add_textbox(Inches(x + 0.22), Inches(y + 0.14), Inches(card_width - 0.38), Inches(card_height - 0.18))
        frame = text_box.text_frame
        frame.word_wrap = True
        frame.clear()
        paragraph = frame.paragraphs[0]
        _add_label_detail_runs(paragraph, item, size=12.5, label_color=theme["accent"])


def _add_cards(slide, items, template="strategy_brief"):
    theme = _ppt_theme(template)
    items = _compact_bullets(items, limit=6, char_limit=120)
    if template == "financial_review":
        positions = [
            (0.92, 2.55), (6.72, 2.55),
            (0.92, 3.48), (6.72, 3.48),
            (0.92, 4.41), (6.72, 4.41),
        ]
        card_width = 5.45
        card_height = 0.74
        for index, item in enumerate(items[:6]):
            x, y = positions[index]
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(card_width), Inches(card_height))
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(255, 255, 255)
            card.line.color.rgb = RGBColor(*theme["soft"])

            number = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.46), Inches(card_height))
            number.fill.solid()
            number.fill.fore_color.rgb = RGBColor(*(theme["accent"] if index % 2 == 0 else theme["secondary"]))
            number.line.fill.background()

            number_box = slide.shapes.add_textbox(Inches(x), Inches(y + 0.2), Inches(0.46), Inches(0.18))
            number_frame = number_box.text_frame
            number_frame.clear()
            number_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            number_run = number_frame.paragraphs[0].add_run()
            number_run.text = str(index + 1)
            _set_run_style(number_run, size=9, bold=True, color=(255, 255, 255))

            text_box = slide.shapes.add_textbox(Inches(x + 0.62), Inches(y + 0.13), Inches(card_width - 0.8), Inches(card_height - 0.12))
            frame = text_box.text_frame
            frame.word_wrap = True
            frame.clear()
            paragraph = frame.paragraphs[0]
            _add_label_detail_runs(paragraph, item, size=11.5, label_color=theme["accent"])
        return

    if len(items) <= 3:
        card_width = 3.55
        card_height = 1.85
        start_x = 1.18 if len(items) == 3 else 2.95 if len(items) == 2 else 4.9
        positions = [(start_x + index * 3.95, 2.88) for index in range(len(items))]
    elif len(items) == 4:
        card_width = 5.2
        card_height = 1.55
        positions = [(1.05, 2.62), (7.0, 2.62), (1.05, 4.38), (7.0, 4.38)]
    else:
        card_width = 3.78
        card_height = 1.58
        positions = [
            (0.78, 2.58), (4.78, 2.58), (8.78, 2.58),
            (0.78, 4.32), (4.78, 4.32), (8.78, 4.32),
        ]
    icon_cycle = ["☁", "▣", "👥", "◆", "↗", "★"]
    for index, item in enumerate(items[:6]):
        x, y = positions[index]
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(card_width), Inches(card_height))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(255, 255, 255)
        card.line.color.rgb = RGBColor(*theme["soft"] if template == "consulting_memo" else PPT_LINE)

        number = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.18), Inches(y + 0.16), Inches(0.33), Inches(0.33))
        number.fill.solid()
        number.fill.fore_color.rgb = RGBColor(*(theme["accent"] if index < 3 else theme["secondary"]))
        number.line.fill.background()

        number_box = slide.shapes.add_textbox(Inches(x + 0.18), Inches(y + 0.205), Inches(0.33), Inches(0.16))
        number_frame = number_box.text_frame
        number_frame.clear()
        number_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        number_run = number_frame.paragraphs[0].add_run()
        number_run.text = icon_cycle[index % len(icon_cycle)]
        _set_run_style(number_run, size=8.5, bold=True, color=(255, 255, 255))

        text_box = slide.shapes.add_textbox(Inches(x + 0.62), Inches(y + 0.14), Inches(card_width - 0.82), Inches(card_height - 0.24))
        frame = text_box.text_frame
        frame.word_wrap = True
        frame.margin_left = Inches(0.02)
        frame.margin_right = Inches(0.02)
        frame.clear()
        paragraph = frame.paragraphs[0]
        _add_label_detail_runs(paragraph, item, size=13.5, label_color=theme["accent"])


def _add_swot_cards(slide, groups):
    labels = [
        ("Strengths", PPT_GREEN, "S"),
        ("Weaknesses", PPT_WARN, "W"),
        ("Opportunities", PPT_BLUE, "O"),
        ("Threats", (129, 94, 38), "T"),
    ]
    positions = [(0.85, 2.5), (6.85, 2.5), (0.85, 4.35), (6.85, 4.35)]
    card_width = 5.55
    card_height = 1.45

    for index, (label, color, initial) in enumerate(labels):
        x, y = positions[index]
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(card_width), Inches(card_height))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(255, 255, 255)
        card.line.color.rgb = RGBColor(*PPT_LINE)

        badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.18), Inches(y + 0.18), Inches(0.38), Inches(0.38))
        badge.fill.solid()
        badge.fill.fore_color.rgb = RGBColor(*color)
        badge.line.fill.background()

        badge_text = slide.shapes.add_textbox(Inches(x + 0.18), Inches(y + 0.235), Inches(0.38), Inches(0.16))
        badge_frame = badge_text.text_frame
        badge_frame.clear()
        badge_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        badge_run = badge_frame.paragraphs[0].add_run()
        badge_run.text = initial
        _set_run_style(badge_run, size=8.5, bold=True, color=(255, 255, 255))

        heading = slide.shapes.add_textbox(Inches(x + 0.68), Inches(y + 0.16), Inches(1.9), Inches(0.32))
        heading_frame = heading.text_frame
        heading_frame.clear()
        heading_run = heading_frame.paragraphs[0].add_run()
        heading_run.text = label
        _set_run_style(heading_run, size=14.5, bold=True, color=color)

        points = [_shorten_text(item, 72) for item in groups.get(label, [])[:2]]
        body = "; ".join(point for point in points if point) or "No major item identified."
        body_box = slide.shapes.add_textbox(Inches(x + 0.68), Inches(y + 0.55), Inches(card_width - 0.95), Inches(0.58))
        body_frame = body_box.text_frame
        body_frame.word_wrap = True
        body_frame.clear()
        body_run = body_frame.paragraphs[0].add_run()
        body_run.text = body
        _set_run_style(body_run, size=11.5, color=PPT_INK)


def _add_kpi_cards(slide, metrics, template="strategy_brief"):
    theme = _ppt_theme(template)
    metrics = _compact_bullets(metrics, limit=5, char_limit=95)
    positions = [
        (0.78, 2.58), (3.25, 2.58), (5.72, 2.58), (8.19, 2.58), (10.66, 2.58),
    ]
    card_width = 2.1
    card_height = 2.65

    for index, metric in enumerate(metrics[:5]):
        x, y = positions[index]
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(card_width), Inches(card_height))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(255, 255, 255)
        card.line.color.rgb = RGBColor(*PPT_LINE)

        badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.18), Inches(y + 0.18), Inches(0.46), Inches(0.46))
        badge.fill.solid()
        badge.fill.fore_color.rgb = RGBColor(*theme["accent"])
        badge.line.fill.background()

        badge_text = slide.shapes.add_textbox(Inches(x + 0.18), Inches(y + 0.25), Inches(0.46), Inches(0.2))
        badge_frame = badge_text.text_frame
        badge_frame.clear()
        badge_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        badge_run = badge_frame.paragraphs[0].add_run()
        badge_run.text = "▣"
        _set_run_style(badge_run, size=11, bold=True, color=(255, 255, 255))

        label, detail = _split_label(metric)
        title = label or f"KPI {index + 1}"
        body = detail if label else metric

        title_box = slide.shapes.add_textbox(Inches(x + 0.22), Inches(y + 0.82), Inches(card_width - 0.44), Inches(0.48))
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        title_frame.clear()
        title_run = title_frame.paragraphs[0].add_run()
        title_run.text = _shorten_text(title, 32)
        _set_run_style(title_run, size=15.5, bold=True, color=theme["navy"])

        body_box = slide.shapes.add_textbox(Inches(x + 0.22), Inches(y + 1.32), Inches(card_width - 0.44), Inches(0.95))
        body_frame = body_box.text_frame
        body_frame.word_wrap = True
        body_frame.clear()
        body_run = body_frame.paragraphs[0].add_run()
        body_run.text = _shorten_text(body, 82)
        _set_run_style(body_run, size=12, color=PPT_INK)


def _risk_pair(text):
    label, detail = _split_label(text)
    if label and re.search(r"mitigat|reduce|address|control|focus|through|by ", detail, re.I):
        parts = re.split(r"\b(?:mitigate|mitigation|reduce|address|control|through|by)\b[:\s-]*", detail, maxsplit=1, flags=re.I)
        if len(parts) == 2 and parts[1].strip():
            return f"{label}: {parts[0].strip(' .;-')}", parts[1].strip(" .;-")
    if ";" in text:
        risk, mitigation = text.split(";", 1)
        return risk.strip(" .;-"), mitigation.strip(" .;-")
    if " - " in text:
        risk, mitigation = text.split(" - ", 1)
        return risk.strip(" .;-"), mitigation.strip(" .;-")
    return text, "Define owner, monitoring signal, and mitigation action."


def _add_risk_mitigation(slide, risks, template="strategy_brief"):
    risks = _compact_bullets(risks, limit=4, char_limit=150)
    headers = ["Risk", "Mitigation"]
    rows = [_risk_pair(risk) for risk in risks]
    table_data = {
        "headers": headers,
        "rows": [[_shorten_text(risk, 92), _shorten_text(mitigation, 105)] for risk, mitigation in rows],
    }
    if template == "consulting_memo":
        _add_table_cards(slide, table_data, template)
    else:
        _add_table_slide(slide, table_data, template)


def _table_shape_config(headers, rows):
    col_count = len(headers)
    financial_headers = {"metric", "latest value", "period / source", "business meaning"}
    is_financial_table = financial_headers.issubset({header.lower() for header in headers})
    if is_financial_table:
        column_widths = [1.75, 2.05, 2.45, 5.8]
        cell_limits = [28, 34, 34, 68]
    elif col_count == 2:
        column_widths = [3.15, 8.9]
        cell_limits = [62, 110]
    elif col_count == 3:
        column_widths = [2.55, 4.35, 5.15]
        cell_limits = [42, 72, 82]
    elif col_count == 4:
        column_widths = [1.75, 2.45, 3.15, 4.7]
        cell_limits = [30, 42, 58, 68]
    else:
        column_widths = [12.05 / col_count] * col_count
        cell_limits = [70] * col_count
    return col_count, is_financial_table, column_widths, cell_limits


def _add_table_cards(slide, table_data, template="consulting_memo"):
    theme = _ppt_theme(template)
    headers = table_data.get("headers") or ["Segment", "Offering", "Customer Value", "Strategic Importance"]
    col_count, is_financial_table, _, cell_limits = _table_shape_config(headers, table_data.get("rows") or [])
    max_rows = 5 if is_financial_table or col_count >= 4 else 6
    rows = (table_data.get("rows") or [])[:max_rows]
    if not rows:
        rows = _fallback_analysis_table([], [], [], [])["rows"][:max_rows]

    start_y = 2.48
    row_gap = 0.12
    card_height = min(0.74, (4.2 - (len(rows) - 1) * row_gap) / max(1, len(rows)))
    for row_index, row in enumerate(rows):
        y = start_y + row_index * (card_height + row_gap)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.82), Inches(y), Inches(11.75), Inches(card_height))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(255, 255, 255)
        card.line.color.rgb = RGBColor(*theme["soft"])

        stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.82), Inches(y), Inches(0.12), Inches(card_height))
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = RGBColor(*(theme["accent"] if row_index % 2 == 0 else theme["secondary"]))
        stripe.line.fill.background()

        title = _shorten_table_text(row[0] if row else "", 34)
        title_box = slide.shapes.add_textbox(Inches(1.08), Inches(y + 0.13), Inches(2.05), Inches(card_height - 0.16))
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        title_frame.clear()
        title_run = title_frame.paragraphs[0].add_run()
        title_run.text = title
        _set_run_style(title_run, size=11.5, bold=True, color=theme["accent"])

        details = []
        for col_index in range(1, min(col_count, len(row))):
            header = headers[col_index] if col_index < len(headers) else ""
            value = _shorten_table_text(row[col_index], cell_limits[col_index] if col_index < len(cell_limits) else 56)
            if value:
                details.append(f"{header}: {value}")
        detail_text = "   |   ".join(details[:3])
        detail_box = slide.shapes.add_textbox(Inches(3.35), Inches(y + 0.12), Inches(8.85), Inches(card_height - 0.14))
        detail_frame = detail_box.text_frame
        detail_frame.word_wrap = True
        detail_frame.clear()
        detail_run = detail_frame.paragraphs[0].add_run()
        detail_run.text = detail_text
        _set_run_style(detail_run, size=10.5, color=PPT_INK)


def _add_table_slide(slide, table_data, template="strategy_brief"):
    theme = _ppt_theme(template)
    headers = table_data.get("headers") or ["Segment", "Offering", "Customer Value", "Strategic Importance"]
    col_count, is_financial_table, column_widths, cell_limits = _table_shape_config(headers, table_data.get("rows") or [])
    max_rows = 5 if is_financial_table or col_count >= 4 else 6
    rows = (table_data.get("rows") or [])[:max_rows]

    if not rows:
        rows = _fallback_analysis_table([], [], [], [])["rows"]

    row_count = len(rows) + 1
    table_top = 2.48
    table_height = 4.25
    table_shape = slide.shapes.add_table(row_count, col_count, Inches(0.62), Inches(table_top), Inches(12.05), Inches(table_height))
    table = table_shape.table

    table_font_size = 9.2 if col_count >= 4 else 10.5
    for index, width in enumerate(column_widths[:col_count]):
        table.columns[index].width = Inches(width)

    header_height = 0.54
    body_height = max(0.62, (table_height - header_height) / max(1, len(rows)))
    table.rows[0].height = Inches(header_height)
    for row_index in range(1, row_count):
        table.rows[row_index].height = Inches(body_height)

    for col_index, header in enumerate(headers):
        cell = table.cell(0, col_index)
        cell.text = header
        cell.margin_left = Inches(0.07)
        cell.margin_right = Inches(0.07)
        cell.margin_top = Inches(0.04)
        cell.margin_bottom = Inches(0.04)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(*theme["navy"])
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        for paragraph in cell.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                _set_run_style(run, size=10.8, bold=True, color=(255, 255, 255))

    for row_index, row in enumerate(rows, start=1):
        for col_index in range(col_count):
            value = row[col_index] if col_index < len(row) else ""
            cell = table.cell(row_index, col_index)
            cell.text = _shorten_table_text(value, cell_limits[col_index] if col_index < len(cell_limits) else 70)
            cell.margin_left = Inches(0.07)
            cell.margin_right = Inches(0.07)
            cell.margin_top = Inches(0.03)
            cell.margin_bottom = Inches(0.03)
            cell.fill.solid()
            if template == "financial_review":
                cell.fill.fore_color.rgb = RGBColor(255, 255, 255) if row_index % 2 else RGBColor(*theme["soft"])
            else:
                cell.fill.fore_color.rgb = RGBColor(255, 255, 255) if row_index % 2 else RGBColor(242, 247, 252)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.LEFT
                for run in paragraph.runs:
                    _set_run_style(run, size=table_font_size, bold=(col_index == 0), color=theme["accent"] if col_index == 0 else PPT_INK)


def _add_footer(slide, slide_number):
    line = slide.shapes.add_shape(
        1,
        Inches(0.65),
        Inches(6.85),
        Inches(12.05),
        Inches(0.02),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(*PPT_LINE)
    line.line.color.rgb = RGBColor(*PPT_LINE)

    footer = slide.shapes.add_textbox(Inches(0.7), Inches(6.95), Inches(11.8), Inches(0.25))
    footer_frame = footer.text_frame
    footer_frame.clear()
    paragraph = footer_frame.paragraphs[0]
    paragraph.alignment = 2
    run = paragraph.add_run()
    run.text = f"Agentic Document Intelligence System | {slide_number}"
    _set_run_style(run, size=8.5, color=PPT_MUTED)


def _slide_matches(slide_data, *terms):
    haystack = " ".join([
        slide_data.get("title", ""),
        slide_data.get("subtitle", ""),
        slide_data.get("kicker", ""),
    ]).lower()
    return any(term in haystack for term in terms)


def _first_slide(slides, *terms):
    for slide_data in slides:
        if _slide_matches(slide_data, *terms):
            return slide_data
    return None


def _dedupe_slides(slides, limit=6):
    selected = []
    seen = set()
    for slide_data in slides:
        if not slide_data:
            continue
        key = slide_data.get("title", "")
        if key in seen:
            continue
        selected.append(slide_data)
        seen.add(key)
        if len(selected) >= limit:
            break
    return selected


def _fallback_content_slides(sections):
    summary_slide = {
        "kind": "bullets",
        "title": "Executive Summary Highlights The Main Decisions",
        "subtitle": "Opportunity, positioning, risk, and recommendation",
        "bullets": sections["summary"],
        "kicker": "summary",
    }
    position_slide = {
        "kind": "bullets",
        "title": "Market Position Shows Where The Offer Can Win",
        "subtitle": "Target segments, differentiation, and value logic",
        "bullets": sections["market_position"],
        "kicker": "position",
    }
    portfolio_slide = {
        "kind": "cards",
        "title": "Portfolio Snapshot Organizes The Core Offerings",
        "subtitle": "Primary services, capabilities, and market components",
        "items": sections["offerings"],
        "kicker": "portfolio",
    }
    comparison_slide = {
        "kind": "bullets",
        "title": "Comparative Analysis Clarifies Strengths And Gaps",
        "subtitle": "Competitive context, trade-offs, and buying criteria",
        "bullets": sections["comparative"],
        "kicker": "comparison",
    }
    strategy_slide = {
        "kind": "bullets",
        "title": "Strategic Moves Translate Analysis Into Action",
        "subtitle": "Recommended actions and expected business impact",
        "bullets": sections["strategy"],
        "kicker": "strategy",
    }
    table_slide = {
        "kind": "table",
        "title": "Segment Table Connects Offerings To Value",
        "subtitle": "Structured comparison across customers, offerings, and strategic importance",
        "table": sections["table"],
        "kicker": "table",
    }
    risk_slide = {
        "kind": "risk_matrix",
        "title": "Key Risks Need Clear Mitigation Paths",
        "subtitle": "Key risks and practical controls",
        "bullets": sections["risks"],
        "kicker": "risk",
    }
    metric_slide = {
        "kind": "kpi",
        "title": "Success Metrics Show Whether The Portfolio Works",
        "subtitle": "Indicators to track portfolio performance",
        "bullets": sections["metrics"],
        "kicker": "metrics",
    }

    return [summary_slide, position_slide, portfolio_slide, comparison_slide, strategy_slide, table_slide]


def _template_content_slides(sections):
    extra_slides = sections.get("extra_slides") or []
    if not extra_slides:
        return _fallback_content_slides(sections)

    company = _first_slide(extra_slides, "company snapshot", "overview")
    financial = _first_slide(extra_slides, "financial", "revenue", "profitability")
    portfolio = _first_slide(extra_slides, "product portfolio", "portfolio")
    market = _first_slide(extra_slides, "target market", "reach")
    competition = _first_slide(extra_slides, "competition", "positioning")
    growth = _first_slide(extra_slides, "growth", "strategic path")
    risk = _first_slide(extra_slides, "risk")
    metrics = _first_slide(extra_slides, "metrics")

    candidates = [company, financial, portfolio, market, competition, growth or risk or metrics]
    remaining = [slide for slide in extra_slides if slide not in candidates]
    return _dedupe_slides([*candidates, *remaining], limit=6)


def _build_ppt_slide_plan(sections, reference_labels, template):
    template = _normalized_ppt_template(template)
    template_names = {
        "strategy_brief": "Classic Consulting Layout",
        "financial_review": "Split-Column Layout",
        "consulting_memo": "Card-Based Layout",
    }
    content_slides = _template_content_slides(sections)[:6]
    return [
        {
            "kind": "cover",
            "title": sections["title"],
            "subtitle": template_names[template],
            "body": f"Prepared by Agentic Document Intelligence System\n{datetime.now().strftime('%d %b %Y, %I:%M %p')}",
        },
        *content_slides,
        {
            "kind": "bullets",
            "title": "Reference Sources",
            "subtitle": "Kept separate from presentation content",
            "bullets": reference_labels or ["No source references were returned for this response."],
            "kicker": "sources",
        },
    ][:8]


def _write_pptx_report(file_path, query, answer, sources, template="strategy_brief"):
    template = _normalized_ppt_template(template)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    sections = _build_ppt_sections(query, answer)
    reference_labels = [_source_label(source) for source in sources]
    slides = _build_ppt_slide_plan(sections, reference_labels, template)

    for index, slide_data in enumerate(slides, start=1):
        slide = prs.slides.add_slide(blank_layout)
        if slide_data["kind"] == "cover":
            _add_cover_slide(slide, slide_data["title"], slide_data["subtitle"], slide_data["body"], template)
            continue

        _apply_slide_background(slide, template)
        _add_title(slide, slide_data["title"], slide_data.get("subtitle"), slide_data.get("kicker", "analysis"), template)
        if slide_data["kind"] == "cards":
            _add_cards(slide, slide_data["items"], template)
        elif slide_data["kind"] == "swot":
            _add_swot_cards(slide, slide_data["swot"])
        elif slide_data["kind"] == "kpi":
            _add_kpi_cards(slide, slide_data["bullets"], template)
        elif slide_data["kind"] == "risk_matrix":
            _add_risk_mitigation(slide, slide_data["bullets"], template)
        elif slide_data["kind"] == "table":
            if template == "consulting_memo":
                _add_table_cards(slide, slide_data["table"], template)
            else:
                _add_table_slide(slide, slide_data["table"], template)
        elif "bullets" in slide_data:
            if slide_data.get("kicker") == "sources":
                _add_bullets(slide, slide_data["bullets"])
            elif template == "financial_review":
                _add_bullets_two_column(slide, slide_data["bullets"], template)
            elif template == "consulting_memo":
                _add_bullet_cards(slide, slide_data["bullets"], template)
            else:
                _add_bullets(slide, slide_data["bullets"])
        else:
            _add_body(slide, slide_data["body"])
        _add_footer(slide, index)

    prs.save(file_path)


@app.post("/api/ask")
async def ask_assistant(request: QueryRequest):
    if not request.query.strip():
        raise HTTPException(status_code=400, detail="Please enter a question.")

    model_settings = _normalized_model_settings(request)
    request_started_at = time.perf_counter()

    try:
        graph_started_at = time.perf_counter()
        result = graph_app.invoke({
            "query": request.query.strip(),
            "presentation_mode": request.presentation_mode,
            **model_settings,
        })
        graph_seconds = round(time.perf_counter() - graph_started_at, 4)
    except Exception as exc:
        raise HTTPException(status_code=500, detail=f"AI engine error: {exc}") from exc

    answer = request.custom_answer.strip() if request.custom_answer else result.get("response", "")
    sources = result.get("sources", [])
    req_format = request.output_format.lower()
    timings = dict(result.get("timings", {}))
    timings["graph_total"] = graph_seconds

    if req_format == "json":
        timings["request_total"] = round(time.perf_counter() - request_started_at, 4)
        return {
            "status": "success",
            "query": request.query,
            "answer": answer,
            "sources": sources,
            "route": result.get("route"),
            "confidence": result.get("confidence", "unknown"),
            "model_settings": model_settings,
            "timings": timings,
        }

    if req_format == "pdf":
        clean_answer = _presentation_text(answer)
        try:
            db_manager.add_generated_knowledge(request.query, clean_answer, sources)
        except Exception as exc:
            print(f"Generated knowledge save skipped: {exc}")

        format_started_at = time.perf_counter()
        file_path = _report_name(request.query, "pdf")
        _write_pdf_report(file_path, request.query, answer, sources)
        timings["format"] = round(time.perf_counter() - format_started_at, 4)
        timings["request_total"] = round(time.perf_counter() - request_started_at, 4)
        return FileResponse(
            str(file_path),
            media_type="application/pdf",
            filename=file_path.name,
            headers=_timing_headers(timings),
        )

    if req_format == "xlsx":
        format_started_at = time.perf_counter()
        file_path = _report_name(request.query, "xlsx")
        _write_xlsx_report(file_path, request.query, answer, sources)
        timings["format"] = round(time.perf_counter() - format_started_at, 4)
        timings["request_total"] = round(time.perf_counter() - request_started_at, 4)
        return FileResponse(
            str(file_path),
            media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            filename=file_path.name,
            headers=_timing_headers(timings),
        )

    if req_format == "pptx":
        format_started_at = time.perf_counter()
        file_path = _report_name(request.query, "pptx")
        clean_answer = _presentation_text(answer)
        try:
            db_manager.add_generated_knowledge(request.query, clean_answer, sources)
        except Exception as exc:
            print(f"Generated knowledge save skipped: {exc}")

        _write_pptx_report(file_path, request.query, answer, sources, request.ppt_template)
        timings["format"] = round(time.perf_counter() - format_started_at, 4)
        timings["request_total"] = round(time.perf_counter() - request_started_at, 4)
        return FileResponse(
            str(file_path),
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            filename=file_path.name,
            headers=_timing_headers(timings),
        )

    raise HTTPException(status_code=400, detail="Invalid format. Choose json, pptx, pdf, or xlsx.")


@app.post("/api/refine")
async def refine_answer(request: RefineRequest):
    if not request.query.strip():
        raise HTTPException(status_code=400, detail="Please enter a question.")
    if not request.current_answer.strip():
        raise HTTPException(status_code=400, detail="No draft answer found to refine.")
    if not request.refinement_prompt.strip():
        raise HTTPException(status_code=400, detail="Please enter a refinement prompt.")

    model_settings = _normalized_model_settings(QueryRequest(
        query=request.query,
        model_provider=request.model_provider,
        model_name=request.model_name,
        temperature=request.temperature,
        max_tokens=request.max_tokens,
        top_p=request.top_p,
    ))

    from agent_graph import _build_llm

    llm = _build_llm(model_settings)
    if llm is None:
        if model_settings["model_provider"] == "openai":
            provider_name = "OpenAI"
            missing_key = "OPENAI_API_KEY"
        elif model_settings["model_provider"] == "groq":
            provider_name = "Groq"
            missing_key = "GROQ_API_KEY or the langchain-groq package"
        else:
            provider_name = "Gemini"
            missing_key = "GOOGLE_API_KEY"
        raise HTTPException(status_code=400, detail=f"{provider_name} is selected, but {missing_key} is not configured.")

    started_at = time.perf_counter()
    prompt = f"""
You are a senior business analyst improving a draft that will later be converted into a professional PowerPoint deck.

Original user request:
{request.query}

Current draft:
{request.current_answer}

Refinement instruction:
{request.refinement_prompt}

Rewrite the draft so it is:
- complete and relevant to the request
- well structured for presentation use
- concise but sufficiently detailed for a mentor/reviewer
- free of repetitive phrasing
- rich in analysis, not only description
- table-friendly where useful, with specific row-level details
- cleanly written without adding a references section inside the body
- business-like and specific enough for mentor review
- fuller in comparative and strategic analysis, not just descriptive
- clear about market logic, target customers, competitive position, risks, and recommendations
- realistic when exact numbers are unavailable; do not invent precise fake statistics
- free of filler phrases like "the provided evidence discusses", "based on limited information", or "I could not find"
- complete enough to generate a PowerPoint directly without another refinement pass
- concise but information-dense; each bullet should contain a specific claim and a business implication

Use this exact section order and these exact headings:
- Executive Summary
- Market Position
- Portfolio Overview
- Comparative Analysis
- Strategic Recommendations
- Suggested Table
- Risks And Mitigations
- Success Metrics

Quality requirements:
- Executive Summary should have 4 strong bullets covering opportunity, positioning, risk, and recommendation
- Market Position should explain target customers, category, differentiation, competition, and pricing/value logic
- Portfolio Overview should explain what each offering does, who it serves, and why it matters commercially
- Comparative Analysis should include meaningful strengths, gaps, risks, alternatives, and competition context
- Strategic Recommendations should be practical, actionable, and tied to business impact
- Suggested Table must be a markdown table with exactly this header:
  | Segment | Offering | Customer Value | Strategic Importance |
  | --- | --- | --- | --- |
- Suggested Table should have exactly 6 data rows with useful row-level detail, not shallow labels
- every table row must fill all 4 columns
- Customer Value and Strategic Importance must be analytical phrases of 8 to 16 words
- do not put bullet points inside table cells
- Risks And Mitigations should identify the main execution/market risks and mitigation actions
- Success Metrics should list measurable indicators such as adoption, retention, revenue mix, usage depth, or expansion

Return only the improved draft. Do not include commentary about what you changed.
"""
    response = llm.invoke(prompt)
    refined_answer = response.content.strip()
    timings = {
        "refine": round(time.perf_counter() - started_at, 4),
        "request_total": round(time.perf_counter() - started_at, 4),
    }
    return {
        "status": "success",
        "query": request.query,
        "answer": refined_answer,
        "timings": timings,
        "model_settings": model_settings,
    }


@app.post("/api/upload")
async def upload_documents(files: list[UploadFile] = File(...)):
    saved_files = []

    for uploaded_file in files:
        extension = Path(uploaded_file.filename or "").suffix.lower()
        if extension not in SUPPORTED_DOCUMENT_EXTENSIONS:
            raise HTTPException(status_code=400, detail=f"{uploaded_file.filename} is not supported.")

        destination = DATA_DIR / _safe_filename(uploaded_file.filename or f"document{extension}")
        with destination.open("wb") as buffer:
            shutil.copyfileobj(uploaded_file.file, buffer)
        saved_files.append(destination.name)

    return {"status": "success", "files": saved_files}


@app.post("/api/ingest")
async def ingest_documents():
    processor = DocumentProcessor(data_dir=str(DATA_DIR))
    chunks = processor.process_all_documents()

    if not chunks:
        deleted_chunks = db_manager.clear()
        return {
            "status": "empty",
            "message": "No documents found. Existing vector memory was cleared.",
            "chunks": 0,
            "deleted_chunks": deleted_chunks,
            "sources": [],
        }

    try:
        db_manager.sync_database(chunks)
    except Exception as exc:
        raise HTTPException(status_code=500, detail=f"Document ingestion failed: {exc}") from exc

    return {
        "status": "success",
        "message": "Documents indexed successfully.",
        "chunks": len(chunks),
        "sources": db_manager.sources(),
    }


@app.post("/api/clear")
async def clear_knowledge_base():
    deleted_chunks = db_manager.clear()
    deleted_files = _clear_uploaded_documents()
    return {
        "status": "success",
        "message": "Knowledge base reset.",
        "deleted_chunks": deleted_chunks,
        "deleted_files": deleted_files,
        "sources": [],
    }


@app.get("/api/status")
async def system_status():
    files = [file.name for file in _document_files()]
    return {
        "status": "online",
        "documents_in_folder": files,
        "indexed_chunks": db_manager.count(),
        "sources": db_manager.sources(),
        "gemini_configured": bool(os.getenv("GOOGLE_API_KEY")),
        "openai_configured": bool(os.getenv("OPENAI_API_KEY")),
        "groq_configured": bool(os.getenv("GROQ_API_KEY")),
        "web_search_configured": bool(os.getenv("TAVILY_API_KEY")),
    }


@app.get("/api/sources")
async def list_sources():
    return {"sources": db_manager.sources()}


@app.get("/health")
async def health_check():
    return {"message": "Agentic Document Intelligence System is online."}


app.mount("/", StaticFiles(directory=str(STATIC_DIR), html=True), name="static")
